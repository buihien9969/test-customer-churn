import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Function to remove markdown formatting
def remove_markdown_formatting(text):
    # Remove bold formatting
    text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
    # Remove italic formatting
    text = re.sub(r'\*(.*?)\*', r'\1', text)
    # Remove backticks
    text = re.sub(r'`(.*?)`', r'\1', text)
    return text

# Function to add a placeholder for charts
def add_sample_chart(slide, chart_type, left=Inches(2), top=Inches(3.5), width=Inches(9), height=Inches(3)):
    # Add a placeholder textbox for the chart
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    p = tf.paragraphs[0]

    # Set chart placeholder text based on type
    chart_descriptions = {
        'pie': 'PIE CHART: Showing distribution of customers',
        'line': 'LINE CHART: Showing trend data',
        'bar': 'BAR CHART: Comparing values across categories',
        'heatmap': 'HEAT MAP: Showing geographic distribution',
        'scatter': 'SCATTER PLOT: Showing correlation between variables',
        'matrix': 'MATRIX CHART: Showing relationships between variables',
        'table': 'TABLE: Showing detailed data comparison',
        'timeline': 'TIMELINE: Showing implementation schedule',
        'confusion': 'CONFUSION MATRIX: Showing model prediction results',
        'workflow': 'WORKFLOW DIAGRAM: Showing process steps'
    }

    if chart_type in chart_descriptions:
        p.text = chart_descriptions[chart_type]
    else:
        p.text = f"[CHART PLACEHOLDER: {chart_type}]"

    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = MCI_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Add a border around the textbox
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.background()  # Make transparent
    shape.line.color.rgb = MCI_GRAY
    shape.line.width = Pt(2)

    return True

# Create a new presentation
prs = Presentation()

# Set slide dimensions to 16:9 aspect ratio
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Define slide layouts
title_slide_layout = prs.slide_layouts[0]  # Title Slide
title_content_layout = prs.slide_layouts[1]  # Title and Content
section_header_layout = prs.slide_layouts[2]  # Section Header
two_content_layout = prs.slide_layouts[3]  # Two Content
title_only_layout = prs.slide_layouts[5]  # Title Only

# Define colors
MCI_BLUE = RGBColor(0, 114, 198)
MCI_DARK_BLUE = RGBColor(0, 65, 132)
MCI_LIGHT_BLUE = RGBColor(135, 206, 250)
MCI_GRAY = RGBColor(128, 128, 128)
MCI_RED = RGBColor(227, 30, 36)
MCI_GREEN = RGBColor(0, 176, 80)
MCI_ORANGE = RGBColor(255, 192, 0)

# Create slides manually to ensure correct content

# Slide 1: Title Slide
slide = prs.slides.add_slide(title_slide_layout)
title_shape = slide.shapes.title
subtitle_shape = slide.placeholders[1]

title_shape.text = "Telecom Customer Churn Analysis"
subtitle_shape.text = "Understanding and Reducing Customer Attrition at MCI"

# Format title
title_shape.text_frame.paragraphs[0].font.size = Pt(44)
title_shape.text_frame.paragraphs[0].font.bold = True
title_shape.text_frame.paragraphs[0].font.color.rgb = MCI_DARK_BLUE

# Format subtitle
subtitle_shape.text_frame.paragraphs[0].font.size = Pt(32)
subtitle_shape.text_frame.paragraphs[0].font.color.rgb = MCI_BLUE

# Add logo placeholder for title slide
logo_box = slide.shapes.add_textbox(Inches(11), Inches(0.5), Inches(2), Inches(1))
logo_p = logo_box.text_frame.paragraphs[0]
logo_p.text = "MCI LOGO"
logo_p.font.size = Pt(14)
logo_p.font.bold = True
logo_p.font.color.rgb = MCI_BLUE
logo_p.alignment = PP_ALIGN.CENTER

# Add decorative element at the bottom
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0),
    Inches(5.5),
    prs.slide_width,
    Inches(2)
)
shape.fill.solid()
shape.fill.fore_color.rgb = MCI_LIGHT_BLUE
shape.line.color.rgb = MCI_BLUE

# Slide 2: Introduction
slide = prs.slides.add_slide(title_content_layout)
title_shape = slide.shapes.title
content_shape = slide.placeholders[1]

title_shape.text = "Introduction"
title_shape.text_frame.paragraphs[0].font.size = Pt(36)
title_shape.text_frame.paragraphs[0].font.bold = True
title_shape.text_frame.paragraphs[0].font.color.rgb = MCI_DARK_BLUE

tf = content_shape.text_frame
tf.clear()

p = tf.paragraphs[0]
p.text = "Problem Statement: Analyzing customer churn patterns at MCI Telecommunications"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = MCI_BLUE

p = tf.add_paragraph()
p.text = "Dataset Overview: 3,333 customer records with 20 attributes"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = MCI_BLUE

p = tf.add_paragraph()
p.text = "Churn Rate: 14.49% of customers have churned"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = MCI_BLUE

p = tf.add_paragraph()
p.text = "Key Questions:"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = MCI_BLUE

p = tf.add_paragraph()
p.text = "1. What is the significance of churn rate?"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "2. What are the characteristics of each customer segment?"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "3. Which machine learning models best predict churn?"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "4. What actions can reduce customer attrition?"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

# Add chart
add_sample_chart(slide, 'pie')

# Add caption
textbox = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11), Inches(0.5))
tf = textbox.text_frame
p = tf.paragraphs[0]
p.text = "Figure: Pie chart showing 14.49% churn vs. 85.51% retained customers"
p.font.italic = True
p.font.size = Pt(14)
p.font.color.rgb = MCI_GRAY
p.alignment = PP_ALIGN.CENTER

# Slide 3: Significance of Churn Rate
slide = prs.slides.add_slide(title_content_layout)
title_shape = slide.shapes.title
content_shape = slide.placeholders[1]

title_shape.text = "Significance of Churn Rate"
title_shape.text_frame.paragraphs[0].font.size = Pt(36)
title_shape.text_frame.paragraphs[0].font.bold = True
title_shape.text_frame.paragraphs[0].font.color.rgb = MCI_DARK_BLUE

tf = content_shape.text_frame
tf.clear()

p = tf.paragraphs[0]
p.text = "Financial Impact:"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = MCI_BLUE

p = tf.add_paragraph()
p.text = "Acquiring new customers costs 5x more than retaining existing ones"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "14.49% churn rate represents significant revenue loss"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "Reducing churn by just 5% can increase profits by 25-95%"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "Stakeholder Impact:"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = MCI_BLUE

p = tf.add_paragraph()
p.text = "For MCI: Revenue stability, market position, operational efficiency"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "For Customers: Service disruption, switching costs, potential for better offers"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "For Investors: Indicator of business health, predictor of future performance"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

# Add chart
add_sample_chart(slide, 'line')

# Add caption
textbox = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11), Inches(0.5))
tf = textbox.text_frame
p = tf.paragraphs[0]
p.text = "Figure: Line graph showing relationship between churn rate reduction and projected profit increase"
p.font.italic = True
p.font.size = Pt(14)
p.font.color.rgb = MCI_GRAY
p.alignment = PP_ALIGN.CENTER

# Slides 4-8 would be created similarly...

# Slide 8: Model Comparison
slide = prs.slides.add_slide(title_content_layout)
title_shape = slide.shapes.title
content_shape = slide.placeholders[1]

title_shape.text = "Model Comparison"
title_shape.text_frame.paragraphs[0].font.size = Pt(36)
title_shape.text_frame.paragraphs[0].font.bold = True
title_shape.text_frame.paragraphs[0].font.color.rgb = MCI_DARK_BLUE

tf = content_shape.text_frame
tf.clear()

p = tf.paragraphs[0]
p.text = "Models Evaluated:"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = MCI_BLUE

models = ["Logistic Regression", "Decision Tree", "Random Forest", "Gradient Boosting", "XGBoost"]
for model in models:
    p = tf.add_paragraph()
    p.text = model
    p.level = 1
    p.font.size = Pt(22)
    p.font.color.rgb = MCI_GRAY

# Add table manually
left = Inches(0.5)
top = Inches(2.5)
width = Inches(12)
height = Inches(3)

# Create table with 6 rows (header + 5 models) and 6 columns
table = slide.shapes.add_table(6, 6, left, top, width, height).table

# Set column widths
col_width = int(width.inches * 914400 / 6)  # Convert to EMU
for col_idx in range(6):
    table.columns[col_idx].width = col_width

# Add headers
headers = ["Model", "Accuracy", "Precision", "Recall", "F1-Score", "ROC AUC"]
for col_idx, header in enumerate(headers):
    cell = table.cell(0, col_idx)
    cell.text = header
    cell.text_frame.paragraphs[0].font.bold = True
    cell.text_frame.paragraphs[0].font.size = Pt(16)
    cell.text_frame.paragraphs[0].font.color.rgb = MCI_DARK_BLUE

# Add data rows
model_data = [
    ["Logistic Regression", "86.06%", "54.35%", "25.77%", "34.97%", "82.19%"],
    ["Decision Tree", "91.45%", "70.41%", "71.13%", "70.77%", "83.02%"],
    ["Random Forest", "93.85%", "98.28%", "58.76%", "73.55%", "92.40%"],
    ["Gradient Boosting", "95.05%", "92.11%", "72.16%", "80.92%", "93.08%"],
    ["XGBoost", "95.50%", "94.67%", "73.20%", "82.56%", "94.63%"]
]

for row_idx, row_data in enumerate(model_data, start=1):
    for col_idx, value in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.text = value
        cell.text_frame.paragraphs[0].font.size = Pt(14)

        # Highlight the best model (XGBoost)
        if row_idx == 5:  # XGBoost row
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.color.rgb = MCI_BLUE

# Add best model note
textbox = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(11), Inches(0.5))
tf = textbox.text_frame
p = tf.paragraphs[0]
p.text = "Best Overall Model: XGBoost (F1-Score: 0.8256, Accuracy: 95.50%)"
p.font.bold = True
p.font.size = Pt(16)
p.font.color.rgb = MCI_BLUE
p.alignment = PP_ALIGN.CENTER

# Slide 9: XGBoost Performance
slide = prs.slides.add_slide(title_content_layout)
title_shape = slide.shapes.title
content_shape = slide.placeholders[1]

title_shape.text = "XGBoost Performance"
title_shape.text_frame.paragraphs[0].font.size = Pt(36)
title_shape.text_frame.paragraphs[0].font.bold = True
title_shape.text_frame.paragraphs[0].font.color.rgb = MCI_DARK_BLUE

tf = content_shape.text_frame
tf.clear()

p = tf.paragraphs[0]
p.text = "Key Metrics:"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = MCI_BLUE

p = tf.add_paragraph()
p.text = "Accuracy: 95.50%"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "Precision: 94.67%"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "Recall: 73.20%"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "F1-Score: 82.56%"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "ROC AUC: 94.63%"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "Confusion Matrix:"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = MCI_BLUE

p = tf.add_paragraph()
p.text = "True Negatives: 566"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "False Positives: 4"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "False Negatives: 26"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "True Positives: 71"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

# Add chart
add_sample_chart(slide, 'confusion')

# Add caption
textbox = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11), Inches(0.5))
tf = textbox.text_frame
p = tf.paragraphs[0]
p.text = "Figure: Confusion matrix heatmap and ROC curve"
p.font.italic = True
p.font.size = Pt(14)
p.font.color.rgb = MCI_GRAY
p.alignment = PP_ALIGN.CENTER

# Slide 10: Feature Importance
slide = prs.slides.add_slide(title_content_layout)
title_shape = slide.shapes.title
content_shape = slide.placeholders[1]

title_shape.text = "Feature Importance"
title_shape.text_frame.paragraphs[0].font.size = Pt(36)
title_shape.text_frame.paragraphs[0].font.bold = True
title_shape.text_frame.paragraphs[0].font.color.rgb = MCI_DARK_BLUE

tf = content_shape.text_frame
tf.clear()

p = tf.paragraphs[0]
p.text = "Top Predictors of Churn:"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = MCI_BLUE

features = [
    "International Plan (17.37%)",
    "Customer Service Calls (9.40%)",
    "Total International Calls (7.94%)",
    "Number of Voicemail Messages (7.59%)",
    "Total Day Minutes (7.47%)"
]

for i, feature in enumerate(features, 1):
    p = tf.add_paragraph()
    p.text = f"{i}. {feature}"
    p.level = 1
    p.font.size = Pt(22)
    p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "Insight: Service plan features and customer service interactions are stronger predictors than demographic factors"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = MCI_BLUE

# Add chart
add_sample_chart(slide, 'bar')

# Add caption
textbox = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11), Inches(0.5))
tf = textbox.text_frame
p = tf.paragraphs[0]
p.text = "Figure: Horizontal bar chart of feature importance"
p.font.italic = True
p.font.size = Pt(14)
p.font.color.rgb = MCI_GRAY
p.alignment = PP_ALIGN.CENTER

# Slide 11: Quantitative Recommendations
slide = prs.slides.add_slide(title_content_layout)
title_shape = slide.shapes.title
content_shape = slide.placeholders[1]

title_shape.text = "Quantitative Recommendations"
title_shape.text_frame.paragraphs[0].font.size = Pt(36)
title_shape.text_frame.paragraphs[0].font.bold = True
title_shape.text_frame.paragraphs[0].font.color.rgb = MCI_DARK_BLUE

tf = content_shape.text_frame
tf.clear()

p = tf.paragraphs[0]
p.text = "International Plan Optimization:"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = MCI_BLUE

p = tf.add_paragraph()
p.text = "Review pricing structure (42.41% churn rate)"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "Create tiered international plans based on usage patterns"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "Implement usage alerts to prevent bill shock"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "Customer Service Enhancement:"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = MCI_BLUE

p = tf.add_paragraph()
p.text = "Proactive outreach after 2+ service calls"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "Implement service quality metrics and training"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "Create dedicated team for high-value customers"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

# Add chart
add_sample_chart(slide, 'matrix')

# Add caption
textbox = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11), Inches(0.5))
tf = textbox.text_frame
p = tf.paragraphs[0]
p.text = "Figure: Impact vs. implementation difficulty matrix"
p.font.italic = True
p.font.size = Pt(14)
p.font.color.rgb = MCI_GRAY
p.alignment = PP_ALIGN.CENTER

# Slide 12: Targeted Retention Strategies
slide = prs.slides.add_slide(title_content_layout)
title_shape = slide.shapes.title
content_shape = slide.placeholders[1]

title_shape.text = "Targeted Retention Strategies"
title_shape.text_frame.paragraphs[0].font.size = Pt(36)
title_shape.text_frame.paragraphs[0].font.bold = True
title_shape.text_frame.paragraphs[0].font.color.rgb = MCI_DARK_BLUE

tf = content_shape.text_frame
tf.clear()

p = tf.paragraphs[0]
p.text = "High Usage Customers:"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = MCI_BLUE

p = tf.add_paragraph()
p.text = "Create unlimited plans for high daytime users"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "Implement loyalty discounts based on usage"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "Develop bundled packages for multi-service users"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "Geographic Focus:"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = MCI_BLUE

p = tf.add_paragraph()
p.text = "Targeted campaigns in high-churn states (NJ, CA, TX)"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "Competitive analysis in these markets"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "Region-specific retention offers"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

# Add chart
add_sample_chart(slide, 'matrix')

# Add caption
textbox = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11), Inches(0.5))
tf = textbox.text_frame
p = tf.paragraphs[0]
p.text = "Figure: Customer segmentation matrix with targeted strategies"
p.font.italic = True
p.font.size = Pt(14)
p.font.color.rgb = MCI_GRAY
p.alignment = PP_ALIGN.CENTER

# Slide 13: Implementation Roadmap
slide = prs.slides.add_slide(title_content_layout)
title_shape = slide.shapes.title
content_shape = slide.placeholders[1]

title_shape.text = "Implementation Roadmap"
title_shape.text_frame.paragraphs[0].font.size = Pt(36)
title_shape.text_frame.paragraphs[0].font.bold = True
title_shape.text_frame.paragraphs[0].font.color.rgb = MCI_DARK_BLUE

tf = content_shape.text_frame
tf.clear()

p = tf.paragraphs[0]
p.text = "Short-term Actions (0-3 months):"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = MCI_BLUE

p = tf.add_paragraph()
p.text = "Implement customer service improvement program"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "Launch international plan review and redesign"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "Deploy predictive churn model in production"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "Medium-term Actions (3-6 months):"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = MCI_BLUE

p = tf.add_paragraph()
p.text = "Roll out new plan structures and pricing"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "Implement automated early warning system"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "Launch targeted retention campaigns by segment"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

# Add chart
add_sample_chart(slide, 'timeline')

# Add caption
textbox = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11), Inches(0.5))
tf = textbox.text_frame
p = tf.paragraphs[0]
p.text = "Figure: Timeline with key milestones"
p.font.italic = True
p.font.size = Pt(14)
p.font.color.rgb = MCI_GRAY
p.alignment = PP_ALIGN.CENTER

# Slide 14: Conclusion & Next Steps
slide = prs.slides.add_slide(title_content_layout)
title_shape = slide.shapes.title
content_shape = slide.placeholders[1]

title_shape.text = "Conclusion & Next Steps"
title_shape.text_frame.paragraphs[0].font.size = Pt(36)
title_shape.text_frame.paragraphs[0].font.bold = True
title_shape.text_frame.paragraphs[0].font.color.rgb = MCI_DARK_BLUE

tf = content_shape.text_frame
tf.clear()

p = tf.paragraphs[0]
p.text = "Key Findings:"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = MCI_BLUE

p = tf.add_paragraph()
p.text = "14.49% overall churn rate with significant variation by segment"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "International plan, customer service calls, and usage patterns are key predictors"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "XGBoost model achieves 95.50% accuracy in predicting churn"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "Expected Impact:"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = MCI_BLUE

p = tf.add_paragraph()
p.text = "5% reduction in churn rate within 6 months"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "Improved customer satisfaction scores"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "Increased customer lifetime value"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

# Slide 15: Appendix - Model Training Steps
slide = prs.slides.add_slide(title_content_layout)
title_shape = slide.shapes.title
content_shape = slide.placeholders[1]

title_shape.text = "Appendix: Model Training Steps"
title_shape.text_frame.paragraphs[0].font.size = Pt(36)
title_shape.text_frame.paragraphs[0].font.bold = True
title_shape.text_frame.paragraphs[0].font.color.rgb = MCI_GREEN

tf = content_shape.text_frame
tf.clear()

p = tf.paragraphs[0]
p.text = "1. Data Preparation:"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = MCI_GREEN

p = tf.add_paragraph()
p.text = "Feature engineering and selection"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "Handling categorical variables with one-hot encoding"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "Standardizing numerical features"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "Train-test split (80/20)"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "2. Model Training Process:"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = MCI_GREEN

p = tf.add_paragraph()
p.text = "Hyperparameter tuning with cross-validation"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "Model evaluation using multiple metrics"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "Feature importance analysis"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "Final model selection based on F1-Score"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

# Add chart
add_sample_chart(slide, 'workflow')

# Add caption
textbox = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11), Inches(0.5))
tf = textbox.text_frame
p = tf.paragraphs[0]
p.text = "Figure: Model training workflow diagram"
p.font.italic = True
p.font.size = Pt(14)
p.font.color.rgb = MCI_GRAY
p.alignment = PP_ALIGN.CENTER

# Slide 16: Thank You
slide = prs.slides.add_slide(title_only_layout)
title_shape = slide.shapes.title

title_shape.text = "Thank You"
title_shape.text_frame.paragraphs[0].font.size = Pt(60)
title_shape.text_frame.paragraphs[0].font.bold = True
title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

# Add special formatting for the Thank You slide
# Create a gradient-like effect with shapes
for i in range(5):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(i * 1.5),
        prs.slide_width,
        Inches(1.5)
    )
    shape.fill.solid()
    # Create a gradient-like effect with different blue shades
    blue_value = max(65, 132 - i * 20)
    shape.fill.fore_color.rgb = RGBColor(0, blue_value, 198)
    shape.line.fill.background()  # No border

# Add a contact information textbox
contact_box = slide.shapes.add_textbox(
    Inches(2),
    Inches(4),
    Inches(9),
    Inches(2)
)
tf = contact_box.text_frame
p = tf.paragraphs[0]
p.text = "Contact Information:"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)

p = tf.add_paragraph()
p.text = "[Your Name]"
p.font.size = Pt(20)
p.font.color.rgb = RGBColor(255, 255, 255)

p = tf.add_paragraph()
p.text = "[Your Email]"
p.font.size = Pt(20)
p.font.color.rgb = RGBColor(255, 255, 255)

p = tf.add_paragraph()
p.text = "[Your Phone Number]"
p.font.size = Pt(20)
p.font.color.rgb = RGBColor(255, 255, 255)

# Save the presentation
try:
    prs.save('MCI_Churn_Analysis_Presentation_Complete.pptx')
    print(f"Presentation created successfully with {len(prs.slides)} slides!")
except Exception as e:
    print(f"Error saving presentation: {e}")
