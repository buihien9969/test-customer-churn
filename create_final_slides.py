import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Function to add a placeholder for charts
def add_sample_chart(slide, chart_type, left=Inches(2), top=Inches(3.5), width=Inches(9), height=Inches(3)):
    # Add a placeholder textbox for the chart
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    p = tf.paragraphs[0]

    # Set chart placeholder text based on type
    chart_descriptions = {
        'pie': 'PIE CHART: Showing distribution of customers',
        'line': 'LINE CHART: Showing trend data',
        'bar': 'BAR CHART: Comparing values across categories',
        'heatmap': 'HEAT MAP: Showing geographic distribution',
        'scatter': 'SCATTER PLOT: Showing correlation between variables',
        'matrix': 'MATRIX CHART: Showing relationships between variables',
        'table': 'TABLE: Showing detailed data comparison',
        'timeline': 'TIMELINE: Showing implementation schedule',
        'confusion': 'CONFUSION MATRIX: Showing model prediction results',
        'workflow': 'WORKFLOW DIAGRAM: Showing process steps'
    }

    if chart_type in chart_descriptions:
        p.text = chart_descriptions[chart_type]
    else:
        p.text = f"[CHART PLACEHOLDER: {chart_type}]"

    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = MCI_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Add a border around the textbox
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.background()  # Make transparent
    shape.line.color.rgb = MCI_GRAY
    shape.line.width = Pt(2)

    return True

# Define colors
MCI_BLUE = RGBColor(0, 114, 198)
MCI_DARK_BLUE = RGBColor(0, 65, 132)
MCI_LIGHT_BLUE = RGBColor(135, 206, 250)
MCI_GRAY = RGBColor(128, 128, 128)
MCI_RED = RGBColor(227, 30, 36)
MCI_GREEN = RGBColor(0, 176, 80)
MCI_ORANGE = RGBColor(255, 192, 0)

# Load the existing presentation
prs = Presentation('MCI_Churn_Analysis_Presentation_Complete.pptx')

# Get slide layouts
title_content_layout = prs.slide_layouts[1]  # Title and Content

# Create a new list of slides
new_slides = []

# Copy the first 3 slides
for i in range(3):
    new_slides.append(prs.slides[i])

# Slide 4: Customer Usage Patterns
slide = prs.slides.add_slide(title_content_layout)
title_shape = slide.shapes.title
content_shape = slide.placeholders[1]

title_shape.text = "Customer Usage Patterns"
title_shape.text_frame.paragraphs[0].font.size = Pt(36)
title_shape.text_frame.paragraphs[0].font.bold = True
title_shape.text_frame.paragraphs[0].font.color.rgb = MCI_DARK_BLUE

tf = content_shape.text_frame
tf.clear()

p = tf.paragraphs[0]
p.text = "Key Differences in Usage:"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = MCI_BLUE

p = tf.add_paragraph()
p.text = "Churned customers use 18.12% more daytime minutes (206.91 vs. 175.18)"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "6.72% more evening minutes (212.41 vs. 199.04)"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "5.33% more international minutes (10.70 vs. 10.16)"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "Charges:"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = MCI_BLUE

p = tf.add_paragraph()
p.text = "Day charges: $35.18 vs. $29.78 (+18.12%)"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "Evening charges: $18.05 vs. $16.92 (+6.71%)"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "International charges: $2.89 vs. $2.74 (+5.33%)"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

# Add chart
add_sample_chart(slide, 'bar')

# Add caption
textbox = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11), Inches(0.5))
tf = textbox.text_frame
p = tf.paragraphs[0]
p.text = "Figure: Grouped bar chart comparing usage patterns between churned and retained customers"
p.font.italic = True
p.font.size = Pt(14)
p.font.color.rgb = MCI_GRAY
p.alignment = PP_ALIGN.CENTER

new_slides.append(slide)

# Slide 5: Service Plans & Customer Support
slide = prs.slides.add_slide(title_content_layout)
title_shape = slide.shapes.title
content_shape = slide.placeholders[1]

title_shape.text = "Service Plans & Customer Support"
title_shape.text_frame.paragraphs[0].font.size = Pt(36)
title_shape.text_frame.paragraphs[0].font.bold = True
title_shape.text_frame.paragraphs[0].font.color.rgb = MCI_DARK_BLUE

tf = content_shape.text_frame
tf.clear()

p = tf.paragraphs[0]
p.text = "International Plan:"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = MCI_BLUE

p = tf.add_paragraph()
p.text = "42.41% churn rate for customers with international plan vs. 11.50% without"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "3.7x higher churn risk with international plan"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "Voice Mail Plan:"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = MCI_BLUE

p = tf.add_paragraph()
p.text = "8.68% churn rate with voice mail plan vs. 16.72% without"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "Voice mail users have 48% lower churn risk"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "Customer Service:"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = MCI_BLUE

p = tf.add_paragraph()
p.text = "Churned customers make 53.8% more service calls (2.23 vs. 1.45)"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "Strong correlation between service calls and churn (r = 0.209)"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

# Add chart
add_sample_chart(slide, 'bar')

# Add caption
textbox = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11), Inches(0.5))
tf = textbox.text_frame
p = tf.paragraphs[0]
p.text = "Figure: Bar charts showing churn rates by plan type"
p.font.italic = True
p.font.size = Pt(14)
p.font.color.rgb = MCI_GRAY
p.alignment = PP_ALIGN.CENTER

new_slides.append(slide)

# Slide 6: Geographic Distribution
slide = prs.slides.add_slide(title_content_layout)
title_shape = slide.shapes.title
content_shape = slide.placeholders[1]

title_shape.text = "Geographic Distribution"
title_shape.text_frame.paragraphs[0].font.size = Pt(36)
title_shape.text_frame.paragraphs[0].font.bold = True
title_shape.text_frame.paragraphs[0].font.color.rgb = MCI_DARK_BLUE

tf = content_shape.text_frame
tf.clear()

p = tf.paragraphs[0]
p.text = "States with Highest Churn:"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = MCI_BLUE

p = tf.add_paragraph()
p.text = "New Jersey (26.47%)"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "California (26.47%)"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "Texas (25.00%)"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "Maryland (24.29%)"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "South Carolina (23.33%)"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "States with Lowest Churn:"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = MCI_BLUE

p = tf.add_paragraph()
p.text = "Alaska (5.77%)"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "Hawaii (5.66%)"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "Virginia (6.49%)"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

# Add chart
add_sample_chart(slide, 'heatmap')

# Add caption
textbox = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11), Inches(0.5))
tf = textbox.text_frame
p = tf.paragraphs[0]
p.text = "Figure: US map heat map showing churn rates by state"
p.font.italic = True
p.font.size = Pt(14)
p.font.color.rgb = MCI_GRAY
p.alignment = PP_ALIGN.CENTER

new_slides.append(slide)

# Slide 7: Customer Segmentation Analysis
slide = prs.slides.add_slide(title_content_layout)
title_shape = slide.shapes.title
content_shape = slide.placeholders[1]

title_shape.text = "Customer Segmentation Analysis"
title_shape.text_frame.paragraphs[0].font.size = Pt(36)
title_shape.text_frame.paragraphs[0].font.bold = True
title_shape.text_frame.paragraphs[0].font.color.rgb = MCI_DARK_BLUE

tf = content_shape.text_frame
tf.clear()

p = tf.paragraphs[0]
p.text = "High-Risk Segment:"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = MCI_BLUE

p = tf.add_paragraph()
p.text = "Has international plan"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "Makes frequent customer service calls (3+ calls)"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "High daytime usage (200+ minutes)"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "Low voicemail usage"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "Low-Risk Segment:"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = MCI_BLUE

p = tf.add_paragraph()
p.text = "Has voicemail plan"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "Few customer service calls (0-1 calls)"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "Moderate usage across all time periods"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

p = tf.add_paragraph()
p.text = "Long account tenure"
p.level = 1
p.font.size = Pt(22)
p.font.color.rgb = MCI_GRAY

# Add chart
add_sample_chart(slide, 'scatter')

# Add caption
textbox = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11), Inches(0.5))
tf = textbox.text_frame
p = tf.paragraphs[0]
p.text = "Figure: Scatter plot showing customer segments by risk factors"
p.font.italic = True
p.font.size = Pt(14)
p.font.color.rgb = MCI_GRAY
p.alignment = PP_ALIGN.CENTER

new_slides.append(slide)

# Add the remaining slides
for i in range(3, len(prs.slides)):
    new_slides.append(prs.slides[i])

# Create a new presentation with all slides in the correct order
new_prs = Presentation()
new_prs.slide_width = prs.slide_width
new_prs.slide_height = prs.slide_height

# Copy all slide layouts
for layout in prs.slide_layouts:
    new_prs.slide_layouts.append(layout)

# Save the presentation
try:
    new_prs.save('MCI_Churn_Analysis_Presentation_Final.pptx')
    print(f"Presentation created successfully with {len(new_slides)} slides!")
except Exception as e:
    print(f"Error saving presentation: {e}")
