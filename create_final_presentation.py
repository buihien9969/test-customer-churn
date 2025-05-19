import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Function to remove markdown formatting
def remove_markdown_formatting(text):
    # Remove bold formatting
    text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
    # Remove italic formatting
    text = re.sub(r'\*(.*?)\*', r'\1', text)
    # Remove backticks
    text = re.sub(r'`(.*?)`', r'\1', text)
    return text

# Function to add a placeholder for charts
def add_sample_chart(slide, chart_type, left=Inches(2), top=Inches(3.5), width=Inches(9), height=Inches(3)):
    # Add a placeholder textbox for the chart
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    p = tf.paragraphs[0]

    # Set chart placeholder text based on type
    chart_descriptions = {
        'pie': 'PIE CHART: Showing distribution of customers',
        'line': 'LINE CHART: Showing trend data',
        'bar': 'BAR CHART: Comparing values across categories',
        'heatmap': 'HEAT MAP: Showing geographic distribution',
        'scatter': 'SCATTER PLOT: Showing correlation between variables',
        'matrix': 'MATRIX CHART: Showing relationships between variables',
        'table': 'TABLE: Showing detailed data comparison',
        'timeline': 'TIMELINE: Showing implementation schedule',
        'confusion': 'CONFUSION MATRIX: Showing model prediction results',
        'workflow': 'WORKFLOW DIAGRAM: Showing process steps'
    }

    if chart_type in chart_descriptions:
        p.text = chart_descriptions[chart_type]
    else:
        p.text = f"[CHART PLACEHOLDER: {chart_type}]"

    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = MCI_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Add a border around the textbox
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.background()  # Make transparent
    shape.line.color.rgb = MCI_GRAY
    shape.line.width = Pt(2)

    return True

# Create a new presentation
prs = Presentation()

# Set slide dimensions to 16:9 aspect ratio
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Define slide layouts
title_slide_layout = prs.slide_layouts[0]  # Title Slide
title_content_layout = prs.slide_layouts[1]  # Title and Content
section_header_layout = prs.slide_layouts[2]  # Section Header
two_content_layout = prs.slide_layouts[3]  # Two Content
title_only_layout = prs.slide_layouts[5]  # Title Only

# Define colors
MCI_BLUE = RGBColor(0, 114, 198)
MCI_DARK_BLUE = RGBColor(0, 65, 132)
MCI_LIGHT_BLUE = RGBColor(135, 206, 250)
MCI_GRAY = RGBColor(128, 128, 128)
MCI_RED = RGBColor(227, 30, 36)
MCI_GREEN = RGBColor(0, 176, 80)
MCI_ORANGE = RGBColor(255, 192, 0)

# Read the markdown file
with open('MCI_Churn_Analysis_Presentation_New.md', 'r', encoding='utf-8') as f:
    md_content = f.read()

# Split the markdown into slides more carefully
raw_slides = md_content.split('---')
slides_md = []

for slide in raw_slides:
    slide = slide.strip()
    if slide:  # Only add non-empty slides
        slides_md.append(slide)

print(f"Found {len(slides_md)} slides in the markdown file")

# Process each slide
for i, slide_md in enumerate(slides_md):
    print(f"Processing slide {i+1}: {slide_md[:50]}...")
    
    # Extract title and content
    lines = slide_md.split('\n')
    title = ""
    content = []

    for line in lines:
        if line.startswith('# '):
            title = line[2:].strip()
        elif line.startswith('## '):
            title = line[3:].strip()
        else:
            content.append(line)

    content_text = '\n'.join(content).strip()

    # Create slide based on content
    if i == 0:  # Title slide
        slide = prs.slides.add_slide(title_slide_layout)
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]

        # Extract subtitle if available
        subtitle = ""
        if len(lines) > 1 and lines[1].startswith('## '):
            subtitle = lines[1][3:].strip()
            # Remove the subtitle line from content
            content_text = '\n'.join(content[1:]).strip()

        title_shape.text = title
        subtitle_shape.text = subtitle

        # Format title
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = MCI_DARK_BLUE

        # Format subtitle
        subtitle_shape.text_frame.paragraphs[0].font.size = Pt(32)
        subtitle_shape.text_frame.paragraphs[0].font.color.rgb = MCI_BLUE

        # Add logo placeholder for title slide
        logo_box = slide.shapes.add_textbox(Inches(11), Inches(0.5), Inches(2), Inches(1))
        logo_p = logo_box.text_frame.paragraphs[0]
        logo_p.text = "MCI LOGO"
        logo_p.font.size = Pt(14)
        logo_p.font.bold = True
        logo_p.font.color.rgb = MCI_BLUE
        logo_p.alignment = PP_ALIGN.CENTER

        # Add decorative element at the bottom
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0),
            Inches(5.5),
            prs.slide_width,
            Inches(2)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = MCI_LIGHT_BLUE
        shape.line.color.rgb = MCI_BLUE
    else:
        # Check if it's an Appendix slide
        is_appendix = "Appendix" in title
        
        # Regular content slide
        slide = prs.slides.add_slide(title_content_layout)
        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]

        title_shape.text = title

        # Format title
        title_shape.text_frame.paragraphs[0].font.size = Pt(36)
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        # Use different color for Appendix slides
        if is_appendix:
            title_shape.text_frame.paragraphs[0].font.color.rgb = MCI_GREEN
        else:
            title_shape.text_frame.paragraphs[0].font.color.rgb = MCI_DARK_BLUE

        # Process content with bullet points
        tf = content_shape.text_frame
        tf.clear()  # Clear existing content

        # Special handling for Model Comparison slide
        if "Model Comparison" in title:
            # Add bullet points for models evaluated
            p = tf.paragraphs[0]
            p.text = "Models Evaluated:"
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = MCI_BLUE
            
            models = ["Logistic Regression", "Decision Tree", "Random Forest", "Gradient Boosting", "XGBoost"]
            for model in models:
                p = tf.add_paragraph()
                p.text = model
                p.level = 1
                p.font.size = Pt(22)
                p.font.color.rgb = MCI_GRAY
            
            # Add table manually
            left = Inches(0.5)
            top = Inches(2.5)
            width = Inches(12)
            height = Inches(3)
            
            # Create table with 6 rows (header + 5 models) and 6 columns
            table = slide.shapes.add_table(6, 6, left, top, width, height).table
            
            # Set column widths
            col_width = int(width.inches * 914400 / 6)  # Convert to EMU
            for col_idx in range(6):
                table.columns[col_idx].width = col_width
            
            # Add headers
            headers = ["Model", "Accuracy", "Precision", "Recall", "F1-Score", "ROC AUC"]
            for col_idx, header in enumerate(headers):
                cell = table.cell(0, col_idx)
                cell.text = header
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].font.size = Pt(16)
                cell.text_frame.paragraphs[0].font.color.rgb = MCI_DARK_BLUE
            
            # Add data rows
            model_data = [
                ["Logistic Regression", "86.06%", "54.35%", "25.77%", "34.97%", "82.19%"],
                ["Decision Tree", "91.45%", "70.41%", "71.13%", "70.77%", "83.02%"],
                ["Random Forest", "93.85%", "98.28%", "58.76%", "73.55%", "92.40%"],
                ["Gradient Boosting", "95.05%", "92.11%", "72.16%", "80.92%", "93.08%"],
                ["XGBoost", "95.50%", "94.67%", "73.20%", "82.56%", "94.63%"]
            ]
            
            for row_idx, row_data in enumerate(model_data, start=1):
                for col_idx, value in enumerate(row_data):
                    cell = table.cell(row_idx, col_idx)
                    cell.text = value
                    cell.text_frame.paragraphs[0].font.size = Pt(14)
                    
                    # Highlight the best model (XGBoost)
                    if row_idx == 5:  # XGBoost row
                        cell.text_frame.paragraphs[0].font.bold = True
                        cell.text_frame.paragraphs[0].font.color.rgb = MCI_BLUE
            
            # Add best model note
            textbox = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(11), Inches(0.5))
            tf = textbox.text_frame
            p = tf.paragraphs[0]
            p.text = "Best Overall Model: XGBoost (F1-Score: 0.8256, Accuracy: 95.50%)"
            p.font.bold = True
            p.font.size = Pt(16)
            p.font.color.rgb = MCI_BLUE
            p.alignment = PP_ALIGN.CENTER
            
            continue  # Skip the rest of the processing for this slide
        
        # Process markdown content for other slides
        current_level = 0
        for line in content_text.split('\n'):
            line = line.strip()
            if not line or line.startswith('*[Visualization:'):
                continue
                
            # Skip table lines for now (we'll handle tables separately)
            if '|' in line and ('----|' in line or line.count('|') >= 3):
                continue

            # Determine indentation level
            level = 0
            if line.startswith('- '):
                line = line[2:]
                level = 0
            elif line.startswith('  - '):
                line = line[4:]
                level = 1
            elif line.startswith('    - '):
                line = line[6:]
                level = 2
            elif line.startswith('1. ') or line.startswith('2. '):
                line = line[3:]
                level = 0

            # Add paragraph with appropriate level
            if tf.paragraphs:
                p = tf.add_paragraph()
            else:
                p = tf.paragraphs[0]

            # Remove markdown formatting
            clean_line = remove_markdown_formatting(line)
            p.text = clean_line
            p.level = level

            # Format text
            p.font.size = Pt(24 - level * 2)
            if level == 0:
                p.font.bold = True
                if is_appendix:
                    p.font.color.rgb = MCI_GREEN
                else:
                    p.font.color.rgb = MCI_BLUE
            else:
                p.font.bold = False
                p.font.color.rgb = MCI_GRAY

        # Add visualization based on slide content
        vis_match = re.search(r'\*\[Visualization: (.*?)\]\*', content_text)
        if vis_match:
            vis_text = vis_match.group(1).lower()

            # Determine chart type based on visualization description
            chart_type = 'bar'  # default
            if 'pie' in vis_text:
                chart_type = 'pie'
            elif 'line' in vis_text:
                chart_type = 'line'
            elif 'heat map' in vis_text or 'heatmap' in vis_text:
                chart_type = 'heatmap'
            elif 'scatter' in vis_text:
                chart_type = 'scatter'
            elif 'matrix' in vis_text:
                chart_type = 'matrix'
            elif 'table' in vis_text:
                chart_type = 'table'
            elif 'timeline' in vis_text:
                chart_type = 'timeline'
            elif 'confusion' in vis_text:
                chart_type = 'confusion'
            elif 'workflow' in vis_text:
                chart_type = 'workflow'

            # Add chart image
            add_sample_chart(slide, chart_type)

            # Add caption
            left = Inches(1)
            top = Inches(6.5)
            width = Inches(11)
            height = Inches(0.5)

            textbox = slide.shapes.add_textbox(left, top, width, height)
            tf = textbox.text_frame
            p = tf.paragraphs[0]
            p.text = f"Figure: {vis_match.group(1)}"
            p.font.italic = True
            p.font.size = Pt(14)
            p.font.color.rgb = MCI_GRAY
            p.alignment = PP_ALIGN.CENTER

# Add special formatting for the last slide (Thank You)
if len(prs.slides) > 0:
    last_slide = prs.slides[-1]

    # Add a nice background for the Thank You slide
    # Create a gradient-like effect with shapes
    for i in range(5):
        shape = last_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0),
            Inches(i * 1.5),
            prs.slide_width,
            Inches(1.5)
        )
        shape.fill.solid()
        # Create a gradient-like effect with different blue shades
        blue_value = max(65, 132 - i * 20)
        shape.fill.fore_color.rgb = RGBColor(0, blue_value, 198)
        shape.line.fill.background()  # No border

    # Make sure the title is still visible
    title_shape = last_slide.shapes.title
    title_shape.text_frame.paragraphs[0].font.size = Pt(60)
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_shape.text_frame.paragraphs[0].font.bold = True

    # Add a contact information textbox
    contact_box = last_slide.shapes.add_textbox(
        Inches(2),
        Inches(4),
        Inches(9),
        Inches(2)
    )
    tf = contact_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Contact Information:"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    p = tf.add_paragraph()
    p.text = "[Your Name]"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(255, 255, 255)

    p = tf.add_paragraph()
    p.text = "[Your Email]"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(255, 255, 255)

    p = tf.add_paragraph()
    p.text = "[Your Phone Number]"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(255, 255, 255)

# Save the presentation
try:
    prs.save('MCI_Churn_Analysis_Presentation_Final.pptx')
    print(f"Presentation created successfully with {len(prs.slides)} slides!")
except Exception as e:
    print(f"Error saving presentation: {e}")
